import os,uuid,hashlib
from decimal import Decimal
from concurrent.futures import ThreadPoolExecutor,as_completed
import boto3,pdfplumber,pandas as pd,requests
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

splitter=RecursiveCharacterTextSplitter(chunk_size=1000,chunk_overlap=100)

def extract(path):
    ext=path.lower().split('.')[-1]
    if ext=='pdf':
        with pdfplumber.open(path) as p:return '\n'.join(x.extract_text() or '' for x in p.pages)
    if ext in ('xls','xlsx'):
        df=pd.read_excel(path,engine='openpyxl');return df.astype(str).apply(lambda x:' | '.join(x),axis=1).str.cat(sep='\n')
    if ext=='docx': return '\n'.join(p.text for p in Document(path).paragraphs)
    if ext=='pptx':
        return '\n'.join(shape.text for s in Presentation(path).slides for shape in s.shapes if hasattr(shape,'text'))
    if ext=='txt': return open(path,encoding='utf-8').read()
    return ''

def file_hash(path):
    h=hashlib.sha256();
    with open(path,'rb') as f:
        for chunk in iter(lambda:f.read(1024*1024),b''):h.update(chunk)
    return h.hexdigest()

def process_all_files(folder):
    # Preserve the original BaseOne RAG target and API embedding behaviour.
    table=boto3.resource('dynamodb',region_name=os.environ['AWS_REGION']).Table('BaseOneRAG'); api=os.environ['OPENAI_API_KEY']
    paths=[os.path.join(root,f) for root,_,files in os.walk(folder) for f in files]
    def process(path):
        text=extract(path)
        if not text.strip(): return 0
        fh=file_hash(path); existing=table.query(KeyConditionExpression=boto3.dynamodb.conditions.Key('PK').eq('DOC#BaseOne')).get('Items',[])
        if any(x.get('file_hash')==fh for x in existing): return 0
        chunks=splitter.split_text(text)
        def upload(chunk):
            r=requests.post('https://api.openai.com/v1/embeddings',headers={'Authorization':f'Bearer {api}','Content-Type':'application/json'},json={'input':chunk,'model':'text-embedding-3-small'});r.raise_for_status(); emb=r.json()['data'][0]['embedding'];table.put_item(Item={'PK':'DOC#BaseOne','SK':f'CHUNK#{uuid.uuid4()}','text':chunk,'embedding':[Decimal(str(x)) for x in emb],'source':path,'file_hash':fh})
        with ThreadPoolExecutor(max_workers=8) as ex:list(ex.map(upload,chunks))
        return len(chunks)
    total=0
    with ThreadPoolExecutor(max_workers=8) as ex:
        for f in as_completed([ex.submit(process,p) for p in paths]): total+=f.result()
    return {'files':len(paths),'chunks':total}
